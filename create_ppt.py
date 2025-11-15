import os
from pathlib import Path

# Try to import packages; script expects python-pptx and cairosvg to be installed.
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    import cairosvg
except Exception as e:
    print('Missing packages:', e)
    raise

ROOT = Path(__file__).parent
ASSETS = ROOT / 'assets'
OUT_PPT = ROOT / 'LeafDisease_Presentation.pptx'

# Input files
README = ROOT / 'README.md'
DETAILED = ROOT / 'DETAILED_MODEL_ARCHITECTURE.md'
SVG_DETAILED = ASSETS / 'model_architecture_detailed.svg'
SVG_SIMPLE = ASSETS / 'model_structure.svg'

# Output images
PNG_DETAILED = ASSETS / 'model_architecture_detailed.png'
PNG_SIMPLE = ASSETS / 'model_structure.png'

# Convert SVGs to PNG
print('Converting SVGs to PNG (this may take a moment)...')
if SVG_DETAILED.exists():
    cairosvg.svg2png(url=str(SVG_DETAILED), write_to=str(PNG_DETAILED), output_width=1600)
    print('Wrote', PNG_DETAILED)
else:
    print('Warning: detailed SVG not found:', SVG_DETAILED)

if SVG_SIMPLE.exists():
    cairosvg.svg2png(url=str(SVG_SIMPLE), write_to=str(PNG_SIMPLE), output_width=1400)
    print('Wrote', PNG_SIMPLE)
else:
    print('Warning: simple SVG not found:', SVG_SIMPLE)

# Read text sources for speaker notes
def read_file(path):
    if path.exists():
        return path.read_text(encoding='utf-8')
    return ''

readme_text = read_file(README)
detailed_text = read_file(DETAILED)

# Helper to add a slide with title and body bullets
def add_bullets_slide(prs, title, bullets, notes=None, image_path=None, image_left=Inches(5.5)):
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        p = body.add_paragraph() if i>0 else body.paragraphs[0]
        p.text = b
        p.level = 0
        p.font.size = Pt(16)
    if image_path and Path(image_path).exists():
        slide.shapes.add_picture(str(image_path), image_left, Inches(1.5), width=Inches(4.7))
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes
    return slide

# Build presentation
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Title slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = 'Leaf Disease Detection — Model & Pipeline'
slide.placeholders[1].text = 'Text-based retrieval using SentenceTransformers (MiniLM)\nFine-tuned with SoftmaxLoss for improved retrieval'
slide.notes_slide.notes_text_frame.text = 'Overview: This deck explains the data pipeline, model architecture (Transformer internals), training with SoftmaxLoss, and inference via cosine retrieval.'

# Data & corpus slide
bullets = [
    'Corpus: disease_data.json — structured fields (leaf_symptoms, fruit_effects, description)',
    'Synthetic queries generated via templates (train_finetune.py)',
    'Dataset split: ~90% train / 10% val for quick evaluation'
]
notes = 'Show a sample disease entry and explain how templates create farmer-style queries.'
add_bullets_slide(prs, 'Data & Corpus', bullets, notes, image_path=PNG_SIMPLE)

# Model architecture slide (simple)
bullets = [
    'Base encoder: all-MiniLM-L6-v2 (6-layer Transformer)',
    'Pooling: mean-pooling over token embeddings -> sentence embedding',
    'Training head: small MLP with SoftmaxLoss (classification over disease labels)'
]
notes = 'Explain that fine-tuning with SoftmaxLoss organizes embedding space for retrieval; show diagram.'
add_bullets_slide(prs, 'Model Architecture — high level', bullets, notes, image_path=PNG_DETAILED)

# Transformer internals slide
bullets = [
    'Tokenization -> embeddings + positional encodings',
    'Multi-head self-attention: Q,K,V projections, scaled-dot product, softmax weights',
    'Feed-Forward Network: per-position MLP with GeLU activation',
    'Residual connections + LayerNorm between sub-layers'
]
notes = 'Provide math: Attention(Q,K,V)=softmax(QK^T/√d_k)V; show why multi-head helps.'
add_bullets_slide(prs, 'Transformer Internals (what/how/why)', bullets, notes)

# Training details slide
bullets = [
    'Loss: SoftmaxLoss (cross-entropy over labels) applied to pooled embeddings',
    'Optimizer: AdamW (typical for Transformers); warmup steps to stabilize LR',
    'Hyperparams used: epochs=3, batch=16, warmup_steps=10 (script defaults)'
]
notes = 'Explain tradeoffs: small model for CPU and quick iteration; more data or GPU recommended for production.'
add_bullets_slide(prs, 'Training & Hyperparameters', bullets, notes)

# Inference & retrieval slide
bullets = [
    'Precompute corpus embeddings from disease descriptions and cache in app session',
    'At query time: encode user text -> compute cosine similarity with corpus embeddings',
    'Return top-1 (or top-k) with confidence; disambiguation flow asks targeted questions when needed'
]
notes = 'Mention caching embeddings to speed inference and L2-normalization so cosine reduces to dot product.'
add_bullets_slide(prs, 'Inference & Retrieval', bullets, notes)

# Demo slide (app)
bullets = [
    'Streamlit demo: app.py provides UI for inputting quantitative, visual, and weather info',
    'Displays predicted disease, confidence bar, reasoning, and expandable details',
    'Disambiguation helper compares two diseases and asks targeted yes/no questions'
]
notes = 'Encourage running streamlit run app.py during demo; show live interaction.'
add_bullets_slide(prs, 'Demo — Streamlit App', bullets, notes)

# Limitations & next steps slide
bullets = [
    'Text-only approach; potential extension: multimodal (images + text)',
    'Improve data: more curated disease descriptions and augmented queries',
    'Production: GPU training, larger model, and persistent embedding cache'
]
notes = 'Next steps: add image pipeline, increase training data, or adopt contrastive losses for retrieval.'
add_bullets_slide(prs, 'Limitations & Next Steps', bullets, notes)

# Closing slide
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = 'Thanks — Questions & Contact'
body = slide.shapes.placeholders[1].text_frame
body.clear()
body.add_paragraph().text = 'Prepared from project: LeafDiseaseDetection'
body.add_paragraph().text = 'Files used: README.md, DETAILED_MODEL_ARCHITECTURE.md, assets/*.svg'
slide.notes_slide.notes_text_frame.text = 'Offer to walk through the model diagram step-by-step during presentation.'

# Save
prs.save(str(OUT_PPT))
print('Wrote PPT:', OUT_PPT)
